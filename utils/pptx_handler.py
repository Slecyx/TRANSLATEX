import subprocess
import os
from pptx import Presentation
from .translator import translate_text

def convert_pptx_to_pdf(pptx_path: str, output_pdf_path: str):
    """Converts a PPTX file to PDF using LibreOffice."""
    try:
        # Check if libreoffice is available
        subprocess.run(['libreoffice', '--version'], check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        
        output_dir = os.path.dirname(output_pdf_path)
        subprocess.run(['libreoffice', '--headless', '--convert-to', 'pdf', pptx_path, '--outdir', output_dir], check=True)
        
        # LibreOffice saves with the same basename but .pdf extension
        expected_output = pptx_path.replace('.pptx', '.pdf')
        
        if os.path.exists(expected_output):
            if expected_output != output_pdf_path:
                if os.path.exists(output_pdf_path):
                    os.remove(output_pdf_path)
                os.rename(expected_output, output_pdf_path)
            return output_pdf_path
        return None
    except Exception as e:
        print(f"Error converting PPTX to PDF: {e}")
        return None

def translate_pptx(input_path: str, output_path: str, target_lang: str, progress_callback=None):
    """Translates a PPTX file to the target language with batch processing."""
    try:
        from .translator import batch_translate_texts
        prs = Presentation(input_path)
        
        # 1. Collect all runs that need translation
        runs_to_translate = []
        original_texts = []
        
        total_slides = len(prs.slides)
        
        # Traverse to collect text
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip():
                            runs_to_translate.append(run)
                            original_texts.append(run.text)
        
        # 2. Batch Translate
        if original_texts:
            translated_texts = batch_translate_texts(original_texts, target_lang)
            
            # 3. Apply translations
            for run, new_text in zip(runs_to_translate, translated_texts):
                run.text = new_text
        
        # Update progress (Simulated as we do batch processing)
        if progress_callback:
             # Just signal completion or simple step as we can't easily track per-slide in batch mode
             # Or we can do batch per slide? 
             # Batching all at once is faster but progress bar will jump.
             # Let's batch per slide to keep progress bar smooth.
             pass

        # Re-implementation for smooth progress with slide-level batching
        # Reloading to clear optimization above for better UX
        runs_to_translate.clear() 
        original_texts.clear()
        
        for i, slide in enumerate(prs.slides):
            if progress_callback:
                progress_callback(i + 1, total_slides)

            # Collect text for THIS slide
            slide_runs = []
            slide_texts = []
            
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip():
                            slide_runs.append(run)
                            slide_texts.append(run.text)
            
            # Translate slide content in parallel
            if slide_texts:
                slide_translations = batch_translate_texts(slide_texts, target_lang)
                for run, new_text in zip(slide_runs, slide_translations):
                    run.text = new_text

        prs.save(output_path)
        return output_path
    except Exception as e:
        print(f"PPTX Error: {e}")
        return output_path
